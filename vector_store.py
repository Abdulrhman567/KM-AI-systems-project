import magic
import PyPDF2
import logging
import chromadb
import pandas as pd
from io import BytesIO
from typing import List
from schema import Document
from pptx import Presentation
from salesforce import Salesforce
from docx import Document as Documentx
from datetime import datetime, timedelta
from dateutil.relativedelta import relativedelta
from chromadb.utils import embedding_functions
from llama_index.vector_stores.chroma import ChromaVectorStore
from langchain.text_splitter import RecursiveCharacterTextSplitter


class Vector_Store:
    def __init__(
        self,
        sf: Salesforce = None,
        collection_name: str = "test",
        embedding_function=None,
    ):
        """
        Parameters:
            sf: Salesforce object to connect with salesforce.
            collection_name (str): Name for the chroma collection.
            embedding_function: Function to create Embedding for the tokens.
        """

        chroma_client = chromadb.PersistentClient(path="chromaDB")
        __embedding_function = (
            embedding_function
            if embedding_function
            else embedding_functions.SentenceTransformerEmbeddingFunction(
                model_name="intfloat/multilingual-e5-base"
            )
        )
        self.files_semantic_collection = chroma_client.get_or_create_collection(
            name=collection_name, embedding_function=__embedding_function
        )
        self.files_exact_collection = chroma_client.get_or_create_collection(
            name="exact_collection"
        )
        self.assets_semantic_collection = chroma_client.get_or_create_collection(
            name="assets_collection", embedding_function=__embedding_function
        )
        self.assets_exact_collection = chroma_client.get_or_create_collection(
            name="assets_exact_collection"
        )

        self.sf = sf
        self.records_ids = []

    def FileType(self, file_bytes: str):
        """
        Determine the file type.

        Parameters:
            file_bytes (str): bytes for the file as str.
        """

        mime = magic.Magic(mime=True)
        file_type = mime.from_buffer(file_bytes)
        return file_type

    def bytes_to_string(self, bytes: str):
        """
        handle different file types, from bytes to string.

        Parameters:
            bytes (bytes): The file in bytes.
        """

        fileType = self.FileType(bytes)
        fileText = ""

        # CSV files.
        if fileType == "text/csv":
            fileText = str(bytes)
            return fileText

        # Excel files.
        elif fileType == "application/vnd.ms-excel":
            xls = pd.read_excel(BytesIO(bytes), sheet_name=None)
            csv_data = {}
            for sheet_name, df in xls.items():
                csv_data[sheet_name] = df.to_csv(index=False)
            fileText = str(csv_data)
            return fileText

        # Powerpoint files.
        elif (
            fileType
            == "application/vnd.openxmlformats-officedocument.presentationml.presentation"
        ):
            prs = Presentation(BytesIO(bytes))
            text = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text.append(shape.text)
            fileText = "\n".join(text)
            return fileText

        # Text files.
        elif fileType == "text/plain":
            fileText = bytes.decode("utf-8")
            return fileText

        # Word files.
        elif (
            fileType
            == "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
        ):
            doc = Documentx(BytesIO(bytes))
            fileText = "\n".join([paragraph.text for paragraph in doc.paragraphs])
            return fileText

        # PDF files.
        elif fileType == "application/pdf":
            pdf_reader = PyPDF2.PdfReader(BytesIO(bytes))
            for page_num in range(len(pdf_reader.pages)):
                page = pdf_reader.pages[page_num]
                fileText += page.extract_text()
            return fileText

        # Other files
        else:
            print(f"unsuported filetype: {fileType}")

    def time_filters(self, **kwargs):
        """
        Create time filters for the search.

        Parameters:
            last_day (bool): Search for records created in the last day.
            last_month (bool): Search for records created in the last month.
            last_year (bool): Search for records created in the last year.
        """
        current_date = datetime.now()
        if kwargs["last_day"]:
            last_day_datetime = str((current_date - timedelta(days=1)).date())
            print(f"last_day_datetime: {last_day_datetime}")
            print(f"Type of last_day_datetime: {type(last_day_datetime)}")
            return last_day_datetime
        if kwargs["last_month"]:
            last_month_datetime = "-".join(
                str((current_date - relativedelta(months=1)).date()).split("-")[:2]
            )
            print(f"Type of last_month_datetime: {type(last_month_datetime)}")
            print(f"last_month_datetime: {last_month_datetime}")
            return last_month_datetime
        if kwargs["last_year"]:
            last_year_datetime = str(
                (current_date - relativedelta(years=1)).date()
            ).split("-")[0]
            print(f"last_year_datetime: {last_year_datetime}")
            print(f"Type of last_year_datetime: {type(last_year_datetime)}")
            return last_year_datetime

    def filters(self, **kwargs):
        """
        Create filters for the search.

        parameters:
            kwargs (dict): The filters for the search.
        """
        print(f"Filters: {kwargs} from (filters)")
        filtered_filters = {
            filter_name: {"$in": filter_value.split(", ")}
            for filter_name, filter_value in kwargs.items()
            if filter_value is not None
        }
        if len(filtered_filters) == 0:
            return None

        elif len(filtered_filters) == 1:
            print(f"The type of the filtered_filters: {type(filtered_filters)}")
            return filtered_filters

        filters_list = [
            {filter_name: filter_value}
            for filter_name, filter_value in filtered_filters.items()
        ]
        print(f"filters_list oasijfiwhefoiwef: {filters_list}")
        dense_search_filters = {"$or": filters_list}
        print(f"Dense search filters: {dense_search_filters} from (filters)")
        print(f"Dense search filters type: {dense_search_filters} from (filters)")
        return dense_search_filters

    def combine_files_with_assets(self, result_files):
        """
        Combine assets and files into a single dictionary.

        Parameters:
            assets (List[Dict]): List of assets.
            files (List[Dict]): List of files.
        """
        list_of_ids = []
        assets_files_list = []
        try:
            for idx, file in enumerate(result_files):
                file_asset = self.sf.get_asset_by_Id(file["ContentDocumentId"])
                if file_asset["Id"] not in list_of_ids:
                    asset_dict = {
                        "Asset_Id": file_asset["Id"],
                        "Asset_Title": file_asset["Title"],
                        "Creation_Date": file_asset["CreatedDate"]
                        .split("T")[0]
                        .replace("-", "/"),
                        "Description": file_asset["Description__c"],
                        "Asset_Url": file_asset["URL__c"],
                        "Asset_Files": [file],
                    }
                    assets_files_list.append(asset_dict)
                    list_of_ids.append(file_asset["Id"])
                else:
                    assets_files_list[idx].append(file)
        except Exception as e:
            print(e)
        return assets_files_list

        # for asset in assets:
        #     if asset["Id"] in assets_files_dict:
        #         assets_files_dict[asset["Id"]]["asset_description"] = asset["Description__c"]
        #         assets_files_dict[asset["Id"]]["asset_url"] = asset["URL__c"]
        #     else:
        #         assets_files_dict[asset["Id"]] = {
        #             "asset_id": asset["Id"],
        #             "asset_title": asset["Title"],
        #             "creation_date": asset["CreatedDate"],
        #             "description": asset["Description__c"],
        #             "asset_url": asset["URL__c"],
        #             "asset_files": [],
        #         }
        # return assets_files_dict

    def combine_assets_with_files(self):
        pass

    def init_vector_semantic(
        self,
        metadata: List[str] = ["Id", "Title"],
        chunk_size=1000,
        chunk_overlap=20,
        From: str = "",
        Where: str = "",
        OrderBy: str = "",
        Limit: str = "",
    ):
        """
        Pull records from Salesforce to Chroma vector store with Text Splitter for semantic search.

        Parameters:
            chunk_size (int): The size of each chunk in bytes.
            chunk_overlap (int): The size of the overlap between adjacent chunks.
            From (str): Name of the Salesforce object from which to retrieve records.
            Where (str): SQL WHERE clause to filter records.
            OrderBy (str): SQL ORDER BY clause to sort records.
            Limit (str): SQL LIMIT clause to limit the number of records retrieved.
        """

        self.records_ids = self.sf.get_records_ids(
            From=From, Where=Where, OrderBy=OrderBy, Limit=Limit
        )
        if self.files_semantic_collection.count() == 0:

            documents = []

            for id in self.records_ids:
                # Get text and metadata for the record
                record_bytes = self.sf.get_record_by_id(id)
                if self.bytes_to_string(record_bytes):
                    record_text = self.bytes_to_string(record_bytes)
                else:
                    continue
                record_metadata = self.sf.get_metadata_by_id(
                    id, metadata_fields=metadata, From=From
                )

                # create document
                doc = Document(page_content=record_text, metadata=record_metadata)
                documents.append(doc)

            text_splitter = RecursiveCharacterTextSplitter(
                chunk_size=chunk_size, chunk_overlap=chunk_overlap
            )
            splits = text_splitter.split_documents(documents)
            splits_text = [splits[i].page_content for i in range(len(splits))]
            splits_metadata = [splits[i].metadata for i in range(len(splits))]
            splits_id = [str(i) for i in range(len(splits))]
            self.files_semantic_collection.add(
                ids=splits_id, documents=splits_text, metadatas=splits_metadata
            )

    def init_vector_exact(
        self,
        metadata: List[str] = ["Id", "Title"],
        From: str = "",
        Where: str = "",
        OrderBy: str = "",
        Limit: str = "",
    ):
        """
        Pull records from Salesforce into Chroma vector store without text splitter for exact search.

        Parameters:
            metadata (List[str]): List of metadata fields to retrieve. Default fields are 'Id' and 'Title'.
            From (str): Name of the Salesforce object from which to retrieve records.
            Where (str): SQL WHERE clause to filter records.
            OrderBy (str): SQL ORDER BY clause to sort records.
            Limit (str): SQL LIMIT clause to limit the number of records retrieved.
        """
        self.records_ids = self.sf.get_records_ids(
            From=From, Where=Where, OrderBy=OrderBy, Limit=Limit
        )
        if self.files_exact_collection.count() == 0:
            documents_data_list = []
            documents_metadata_list = []
            for id in self.records_ids:
                record_bytes = self.sf.get_record_by_id(id)
                if self.bytes_to_string(record_bytes):
                    documents_data_list.append(self.bytes_to_string(record_bytes))
                else:
                    continue
                documents_metadata_list.append(
                    self.sf.get_metadata_by_id(id, metadata_fields=metadata, From=From)
                )
            documents_ids_list = [str(i) for i in range(len(documents_data_list))]
            self.files_exact_collection.add(
                ids=documents_ids_list,
                metadatas=documents_metadata_list,
                documents=documents_data_list,
            )

    def init_vector_assets(
        self,
        metadata: List[str] = ["Id", "Title"],
        From: str = "",
        Where: str = "",
        OrderBy: str = "",
        Limit: str = "",
    ):
        self.records_ids = self.sf.get_records_ids(
            From=From, Where=Where, OrderBy=OrderBy, Limit=Limit
        )
        if self.assets_semantic_collection.count() == 0:
            documents_data_list = []
            documents_metadata_list = []
            for id in self.records_ids:
                asset_data, asset_metadata = {}, {}
                asset_data = self.sf.get_metadata_by_id(
                    id, metadata_fields=metadata, From="Knowledge__kav"
                )
                asset_metadata["Id"] = asset_data["Id"]
                asset_metadata["KnowledgeArticleId"] = asset_data["KnowledgeArticleId"]
                asset_title = asset_data["Title"]
                asset_summary = asset_data["Summary"]
                if asset_title or asset_summary:
                    documents_data_list.append(f"{asset_title} {asset_summary}")
                else:
                    continue
                documents_metadata_list.append(asset_metadata)
            self.assets_semantic_collection.add(
                ids=[str(id) for id in range(len(documents_data_list))],
                metadatas=documents_metadata_list,
                documents=documents_data_list,
            )

    def init_vector_assets_exact(
        self,
        metadata: List[str] = ["Id", "Title"],
        From: str = "",
        Where: str = "",
        OrderBy: str = "",
        Limit: str = "",
    ):
        self.records_ids = self.sf.get_records_ids(
            From=From, Where=Where, OrderBy=OrderBy, Limit=Limit
        )
        if self.assets_exact_collection.count() == 0:
            documents_data_list = []
            documents_metadata_list = []
            for id in self.records_ids:
                asset_data, asset_metadata = {}, {}
                asset_data = self.sf.get_metadata_by_id(
                    id, metadata_fields=metadata, From="Knowledge__kav"
                )
                asset_metadata["Id"] = asset_data["Id"]
                asset_metadata["KnowledgeArticleId"] = asset_data["KnowledgeArticleId"]
                asset_title = asset_data["Title"]
                asset_summary = asset_data["Summary"]
                if asset_title or asset_summary:
                    documents_data_list.append(f"{asset_title} {asset_summary}")
                else:
                    continue
                documents_metadata_list.append(asset_metadata)
            self.assets_exact_collection.add(
                ids=[str(id) for id in range(len(documents_data_list))],
                metadatas=documents_metadata_list,
                documents=documents_data_list,
            )

    def files_semantic_search(
        self,
        query: str,
        n_results: int = 10,
        max_distance: float = 1.8,
        filters: dict = None,
        time_filter: str = None,
    ):
        """
        Search the vector store of attachments with a query

        Parameters:
            query (str): Search query.
            n_results (int): Number of retrieved records.
            max_distance (float): Maximum distance between the query and documents.
        """

        matching_docs = self.files_semantic_collection.query(
            query_texts=[query],
            n_results=self.files_semantic_collection.count(),
            where=filters,
        )
        unique_list = []
        try:
            results_list = []
            # print(f"Matching docs: {matching_docs} from (files_semantic_search)")
            for i in range(self.files_semantic_collection.count()):
                print(i)
                if matching_docs["metadatas"][0][i]["Id"] not in unique_list:
                    if matching_docs["distances"][0][i] < max_distance:
                        unique_list.append(matching_docs["metadatas"][0][i]["Id"])
                        matching_docs["metadatas"][0][i]["preview"] = (
                            matching_docs["documents"][0][i] + "..."
                        ).replace("\n", "")
                        matching_docs["metadatas"][0][i][
                            "url"
                        ] = f"http://0.0.0.0:8000/records/download?id={matching_docs['metadatas'][0][i]['Id']}"
                        results_list.append(matching_docs["metadatas"][0][i])
                    if len(unique_list) >= n_results:
                        break
        except Exception as e:
            print(f"Error: {e}")
        if time_filter:
            print(time_filter)
            print(type(time_filter))
            time_filtered_resultes = [
                doc for doc in results_list if time_filter in doc["CreatedDate"]
            ]
            print(time_filtered_resultes)
            return time_filtered_resultes
        return results_list

    def assets_semantic_search(
        self,
        query: str,
        n_results: int = 10,
        max_distance: float = 1.8,
        # filters: dict = None,
        # time_filters: str = None,
    ):
        """
        Search the vector store of assets with a query
        """
        matching_assets = self.assets_semantic_collection.query(
            query_texts=[query],
            n_results=self.assets_semantic_collection.count(),
            # where=filters,
        )
        # print(f"Macthing assets: {matching_assets} from assets_semantic_search")
        # print(f"Fitlers: {filters} from (assets_semantic_search)")
        unique_list = []
        try:
            results_list = []
            for i in range(self.assets_semantic_collection.count()):
                if matching_assets["metadatas"][0][i]["Id"] not in unique_list:
                    if matching_assets["distances"][0][i] < max_distance:
                        unique_list.append(matching_assets["metadatas"][0][i]["Id"])
                        matching_assets["metadatas"][0][i]["preview"] = (
                            matching_assets["documents"][0][i] + "..."
                        ).replace("\n", "")
                        results_list.append(matching_assets["metadatas"][0][i])
                    if len(unique_list) >= n_results:
                        break
        except Exception as e:
            print(f"Error: {e}")
        return results_list

    def files_exact_search(
        self,
        search_query: str,
        n_results: int = 10,
        filters: dict = None,
        time_filter: str = None,
    ):
        matching_docs = self.files_exact_collection.query(
            query_texts=[search_query],
            n_results=self.files_exact_collection.count(),
            where_document={"$contains": search_query},
            where=filters,
        )
        unique_list = []
        try:
            results_list = []
            for i in range(len(matching_docs["ids"][0])):
                if matching_docs["metadatas"][0][i]["Id"] not in unique_list:
                    unique_list.append(matching_docs["metadatas"][0][i]["Id"])
                    matching_docs["metadatas"][0][i]["preview"] = (
                        matching_docs["documents"][0][i][:300] + "..."
                    ).replace("\n", "")
                    matching_docs["metadatas"][0][i][
                        "url"
                    ] = f"http://0.0.0.0:8000/records/download?id={matching_docs['metadatas'][0][i]['Id']}"
                    results_list.append(matching_docs["metadatas"][0][i])
                if len(unique_list) >= n_results:
                    break
        except Exception as e:
            print(e)

        if time_filter:
            print(time_filter)
            print(type(time_filter))
            time_filtered_resultes = [
                doc for doc in results_list if time_filter in doc["CreatedDate"]
            ]
            print(time_filtered_resultes)
            return time_filtered_resultes
        return results_list

    def assets_exact_search(
        self,
        search_query: str,
        n_results: int = 1,
    ):
        matching_assets = self.assets_exact_collection.query(
            query_texts=[search_query],
            n_results=self.assets_exact_collection.count(),
            where_document={"$contains": search_query},
        )
        unique_list = []
        results_list = []
        for i in range(len(matching_assets["ids"][0])):
            if matching_assets["metadatas"][0][i]["Id"] not in unique_list:
                unique_list.append(matching_assets["metadatas"][0][i]["Id"])
                matching_assets["metadatas"][0][i]["preview"] = (
                    matching_assets["documents"][0][i][:100] + "..."
                ).replace("\n", "")
                results_list.append(matching_assets["metadatas"][0][i])
            if len(unique_list) >= n_results:
                break
        return results_list

    def chatbot(self, query: str):
        vector_store = ChromaVectorStore(
            chroma_collection=self.files_semantic_collection
        )

    def add_document(
        self,
        chunk_size: int = 1000,
        chunk_overlap: int = 20,
        From: str = "",
        Where: str = "",
        OrderBy: str = "",
        Limit: str = "",
    ):
        """
        Track new files uploaded to salesforce and pull them.

        Parameters:
            chunk_size (bytes): The size of each chunk.
            chunk_overlap: Chunk overlap size.
            From (str): Salesforce object name.
            Where (str): Where clause.
            OrderBy (str): Order by clause.
            Limit (str): Limit clause.
        """
        try:
            old_ids = self.records_ids.copy()
            self.records_ids = []
            self.records_ids = self.sf.get_records_ids(
                From=From, Where=Where, OrderBy=OrderBy, Limit=Limit
            )
            new_ids = set(self.records_ids) - set(old_ids)

            documents = []
            for id in new_ids:
                record_bytes = self.sf.get_record_by_id(id)
                if self.bytes_to_string(record_bytes):
                    record_text = self.bytes_to_string(record_bytes)
                else:
                    continue
                record_metadata = self.sf.get_metadata_by_id(
                    id,
                    metadata_fields=[
                        "Id",
                        "CreatedById",
                        "CreatedDate",
                        "Title",
                        "FileType",
                    ],
                    From=From,
                )

                # create document
                doc = Document(page_content=record_text, metadata=record_metadata)
                self.files_semantic_collection.add
                documents.append(doc)
            text_splitter = RecursiveCharacterTextSplitter(
                chunk_size=chunk_size, chunk_overlap=chunk_overlap
            )
            splits = text_splitter.split_documents(documents)
            splits_text = [splits[i].page_content for i in range(len(splits))]
            splits_metadata = [splits[i].metadata for i in range(len(splits))]
            splits_id = [
                str(i)
                for i in range(
                    self.files_semantic_collection.count(),
                    (len(splits)) + self.files_semantic_collection.count(),
                    1,
                )
            ]
            self.files_semantic_collection.add(
                ids=splits_id, documents=splits_text, metadatas=splits_metadata
            )
        except Exception as e:
            logging.exception("No new uploaded files on salesforce.")

    def delete_document(
        self, From: str = "", Where: str = "", OrderBy: str = "", Limit: str = ""
    ):
        """
        Delete a document from the vector store when it's deleted from salesforce.

        Parameters:
            From (str): Salesforce object name.
            Where (str): Where clause.
            OrderBy (str): Order by clause.
            Limit (str): Limit clause.
        """
        try:
            old_ids = self.records_ids.copy()
            self.records_ids = []
            self.records_ids = self.sf.get_records_ids(
                From=From, Where=Where, OrderBy=OrderBy, Limit=Limit
            )
            deleted_ids = set(old_ids) - set(self.records_ids)
            for id in deleted_ids:
                self.files_semantic_collection.delete(where={"Id": f"{id}"})

        except Exception as e:
            logging.exception("No new deleted files on salesforce.")
